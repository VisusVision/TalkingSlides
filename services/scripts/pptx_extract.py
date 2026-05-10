# services/scripts/pptx_extract.py
"""
Unified lesson-source extraction helper for AI_ACADEMY.

Supported input formats:
  - PPTX  (.pptx)   — slide images via LibreOffice or python-pptx + Pillow
  - PDF   (.pdf)    — page images via pdftoppm or PyMuPDF
  - DOCX  (.docx)   — page images via LibreOffice or python-docx + Pillow
  - TXT   (.txt)    — text paragraphs rendered as images via Pillow
  - Image (.png/.jpg/.jpeg/.webp/.gif) — single-slide image sources

Public API (compatible with existing callers):
  export_slide_images(path, out_dir, resolution=1920)               -> list[str]
  export_slide_images_with_metadata(path, out_dir, resolution=1920) -> dict
  extract_speaker_notes(path, out_dir)                              -> list[str]

Both functions dispatch on file extension so the pipeline worker can pass
any supported file without changes.
"""

from __future__ import annotations

import io
import logging
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Any, List

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Optional imports
# ---------------------------------------------------------------------------
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception:
    Presentation = None
    MSO_SHAPE_TYPE = None

try:
    from PIL import Image, ImageDraw, ImageFont, ImageOps
except Exception:
    Image = None
    ImageDraw = None
    ImageFont = None
    ImageOps = None

try:
    import fitz as _pymupdf  # PyMuPDF
    _HAVE_PYMUPDF = True
except Exception:
    _pymupdf = None
    _HAVE_PYMUPDF = False

try:
    from docx import Document as _DocxDocument
    _HAVE_DOCX = True
except Exception:
    _DocxDocument = None
    _HAVE_DOCX = False

def _env_int(name: str, default: int) -> int:
    try:
        return int(os.environ.get(name, str(default)))
    except (TypeError, ValueError):
        return default


_LIBREOFFICE_EXPORT_TIMEOUT_SECONDS = _env_int("LIBREOFFICE_EXPORT_TIMEOUT_SECONDS", 120)
_LIBREOFFICE_STDIO_TAIL_CHARS = 1200
_LIBREOFFICE_MAX_LISTED_OUTPUT_FILES = 50
_LIBREOFFICE_CONFIG_ENV_VARS = ("SOFFICE_PATH", "LIBREOFFICE_PATH", "LIBREOFFICE_EXECUTABLE")


class LibreOfficeExportError(RuntimeError):
    """Raised when the LibreOffice conversion step fails with diagnostic metadata."""

    def __init__(self, message: str, *, warning_code: str, details: dict[str, Any] | None = None):
        warnings = ["libreoffice_export_failed"]
        if warning_code and warning_code not in warnings:
            warnings.append(warning_code)
        self.warning_code = warning_code
        self.warnings = warnings
        self.details = dict(details or {})
        self.details.setdefault("warning_code", warning_code)
        super().__init__(message)

# Minimal 1×1 white PNG bytes (used as last-resort stub)
_WHITE_1X1_PNG = (
    b"\x89PNG\r\n\x1a\n"
    b"\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x02"
    b"\x00\x00\x00\x90wS\xde"
    b"\x00\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00\x05\x18\xd4n"
    b"\x00\x00\x00\x00IEND\xaeB`\x82"
)

# Supported file extensions
_IMAGE_SOURCE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".gif"}
_SUPPORTED_EXTS = {".pptx", ".pdf", ".docx", ".txt", *_IMAGE_SOURCE_EXTS}


# ---------------------------------------------------------------------------
# Utilities
# ---------------------------------------------------------------------------

def _ensure_dir(path: str) -> Path:
    p = Path(path)
    p.mkdir(parents=True, exist_ok=True)
    return p


def _safe_replace(src: Path, dst: Path) -> None:
    """Move *src* to *dst*, handling cross-device links (different Docker mounts)."""
    try:
        src.replace(dst)
    except OSError:
        shutil.copy2(src, dst)
        try:
            src.unlink()
        except OSError:
            pass


def _tail_text(value: Any, limit: int = _LIBREOFFICE_STDIO_TAIL_CHARS) -> str:
    if value is None:
        return ""
    if isinstance(value, bytes):
        text = value.decode("utf-8", errors="replace")
    else:
        text = str(value)
    return text[-limit:] if len(text) > limit else text


def _command_text(cmd: List[str]) -> str:
    return subprocess.list2cmdline([str(part) for part in cmd])


def _find_soffice_executable() -> str | None:
    for env_name in _LIBREOFFICE_CONFIG_ENV_VARS:
        configured = os.environ.get(env_name)
        if configured and Path(configured).exists():
            return configured
    for name in ("soffice", "libreoffice"):
        found = shutil.which(name)
        if found:
            return found
    if os.name == "nt":
        for root in (os.environ.get("PROGRAMFILES"), os.environ.get("PROGRAMFILES(X86)")):
            if not root:
                continue
            candidate = Path(root) / "LibreOffice" / "program" / "soffice.exe"
            if candidate.exists():
                return str(candidate)
    return None


def _listed_output_files(out_dir: Path) -> list[str]:
    if not out_dir.exists():
        return []
    files: list[str] = []
    try:
        entries = sorted(out_dir.iterdir(), key=lambda p: p.name.lower())
    except OSError:
        return []
    for path in entries[:_LIBREOFFICE_MAX_LISTED_OUTPUT_FILES]:
        try:
            if path.is_file():
                files.append(f"{path.name} ({path.stat().st_size} bytes)")
            elif path.is_dir():
                files.append(f"{path.name}/")
            else:
                files.append(path.name)
        except OSError:
            files.append(path.name)
    if len(entries) > _LIBREOFFICE_MAX_LISTED_OUTPUT_FILES:
        files.append(f"... {len(entries) - _LIBREOFFICE_MAX_LISTED_OUTPUT_FILES} more")
    return files


def _libreoffice_details(
    *,
    cmd: List[str],
    source_path: Path,
    out_dir: Path,
    expected_pdf: Path,
    cwd: Path,
    staged_input_path: Path | None = None,
    profile_dir: Path | None = None,
    return_code: int | None = None,
    stdout: Any = "",
    stderr: Any = "",
) -> dict[str, Any]:
    details = {
        "command": _command_text(cmd),
        "return_code": return_code,
        "stdout_tail": _tail_text(stdout),
        "stderr_tail": _tail_text(stderr),
        "expected_output_path": str(expected_pdf),
        "actual_output_files": _listed_output_files(out_dir),
        "working_directory": str(cwd),
        "input_path": str(source_path),
        "output_directory": str(out_dir),
    }
    if staged_input_path is not None:
        details["staged_input_path"] = str(staged_input_path)
    if profile_dir is not None:
        details["user_profile_path"] = str(profile_dir)
    return details


def _stage_libreoffice_input(source_path: Path, staging_dir: Path) -> Path:
    staging_dir.mkdir(parents=True, exist_ok=True)
    staged_path = staging_dir / source_path.name
    shutil.copy2(source_path, staged_path)
    return staged_path


def _pdf_output_candidates(out_dir: Path) -> list[Path]:
    if not out_dir.exists():
        return []
    try:
        return [path for path in out_dir.iterdir() if path.is_file() and path.suffix.lower() == ".pdf"]
    except OSError:
        return []


def _select_libreoffice_pdf_output(source_path: Path, out_dir: Path) -> Path | None:
    candidates = _pdf_output_candidates(out_dir)
    valid = [path for path in candidates if path.exists() and path.stat().st_size > 0]
    if not valid:
        return None
    matching = [path for path in valid if path.stem.lower() == source_path.stem.lower()]
    selected_pool = matching or valid
    return max(selected_pool, key=lambda path: path.stat().st_mtime)


def _libreoffice_failure_warnings(exc: Exception) -> List[str]:
    warnings = getattr(exc, "warnings", None)
    if isinstance(warnings, list) and warnings:
        return list(dict.fromkeys(str(item) for item in warnings if str(item or "").strip()))
    return ["libreoffice_export_failed"]


def _libreoffice_failure_details(exc: Exception) -> list[dict[str, Any]]:
    details = getattr(exc, "details", None)
    if isinstance(details, dict) and details:
        return [details]
    return [
        {
            "warning_code": "libreoffice_export_failed",
            "message": _tail_text(str(exc), 500),
        }
    ]


def _run(cmd: List[str], check=True, capture=False, timeout: int | None = None, cwd: str | Path | None = None):
    logger.debug("Running: %s", " ".join(cmd))
    try:
        proc = subprocess.run(
            cmd,
            stdout=subprocess.PIPE if capture else None,
            stderr=subprocess.PIPE if capture else None,
            text=True,
            timeout=timeout,
            cwd=str(cwd) if cwd is not None else None,
        )
    except subprocess.TimeoutExpired as exc:
        raise RuntimeError(
            f"Command timed out: {' '.join(cmd)}\n"
            f"stdout: {_tail_text(exc.stdout)}\nstderr: {_tail_text(exc.stderr)}"
        ) from exc
    if check and proc.returncode != 0:
        raise RuntimeError(
            f"Command failed ({proc.returncode}): {' '.join(cmd)}\n"
            f"stdout: {proc.stdout or ''}\nstderr: {proc.stderr or ''}"
        )
    return proc


def source_render_dependency_report() -> dict:
    return {
        "libreoffice_available": bool(_find_soffice_executable()),
        "pdftoppm_available": bool(shutil.which("pdftoppm")),
        "pymupdf_available": bool(_HAVE_PYMUPDF),
        "python_pptx_available": bool(Presentation is not None),
    }


def source_render_dependency_warnings(source_ext: str | None = None) -> List[str]:
    ext = str(source_ext or "").strip().lower()
    if ext and ext not in {".pptx", ".pdf", ".docx"}:
        return []
    report = source_render_dependency_report()
    warnings: List[str] = []
    if ext in {"", ".pptx", ".docx"} and not report["libreoffice_available"]:
        warnings.append("slide_render_dependency_missing_libreoffice")
    if ext in {"", ".pptx", ".docx", ".pdf"} and not report["pdftoppm_available"]:
        warnings.append("slide_render_dependency_missing_pdftoppm")
    return warnings


def _convert_via_libreoffice_to_pdf(source_path: str, out_dir: Path) -> Path:
    """Convert an office document to PDF in *out_dir* and return the PDF path."""
    source = Path(source_path).resolve()
    out_dir = Path(out_dir).resolve()
    out_dir.mkdir(parents=True, exist_ok=True)
    expected_pdf = out_dir / f"{source.stem}.pdf"
    fallback_cwd = source.parent if source.parent.exists() else Path.cwd()
    soffice = _find_soffice_executable()
    if not soffice:
        cmd = [
            "soffice", "--headless", "--nologo", "--nofirststartwizard",
            "--nolockcheck", "--nodefault", "--convert-to", "pdf",
            "--outdir", str(out_dir), str(source),
        ]
        details = _libreoffice_details(
            cmd=cmd,
            source_path=source,
            out_dir=out_dir,
            expected_pdf=expected_pdf,
            cwd=fallback_cwd,
            stderr="LibreOffice executable was not found.",
        )
        raise LibreOfficeExportError(
            "LibreOffice executable was not found",
            warning_code="libreoffice_export_failed",
            details=details,
        )

    with tempfile.TemporaryDirectory(prefix="lo-profile-") as profile_tmp, tempfile.TemporaryDirectory(prefix="lo-input-") as stage_tmp:
        profile_dir = Path(profile_tmp).resolve()
        profile_dir.mkdir(parents=True, exist_ok=True)
        staging_dir = Path(stage_tmp).resolve()
        staged_source = staging_dir / source.name
        lo_cmd = [
            soffice,
            "--headless",
            "--nologo",
            "--nofirststartwizard",
            "--nolockcheck",
            "--nodefault",
            f"-env:UserInstallation={profile_dir.as_uri()}",
            "--convert-to",
            "pdf",
            "--outdir",
            str(out_dir),
            str(staged_source),
        ]
        try:
            staged_source = _stage_libreoffice_input(source, staging_dir)
        except Exception as exc:
            details = _libreoffice_details(
                cmd=lo_cmd,
                source_path=source,
                staged_input_path=staged_source,
                out_dir=out_dir,
                expected_pdf=expected_pdf,
                cwd=staging_dir,
                profile_dir=profile_dir,
                stderr=str(exc),
            )
            raise LibreOfficeExportError(
                "LibreOffice input staging failed",
                warning_code="libreoffice_input_staging_failed",
                details=details,
            ) from exc

        try:
            proc = subprocess.run(
                lo_cmd,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                timeout=_LIBREOFFICE_EXPORT_TIMEOUT_SECONDS,
                cwd=str(staging_dir),
            )
        except subprocess.TimeoutExpired as exc:
            details = _libreoffice_details(
                cmd=lo_cmd,
                source_path=source,
                staged_input_path=staged_source,
                out_dir=out_dir,
                expected_pdf=expected_pdf,
                cwd=staging_dir,
                profile_dir=profile_dir,
                stdout=exc.stdout,
                stderr=exc.stderr,
            )
            raise LibreOfficeExportError(
                "LibreOffice export timed out",
                warning_code="libreoffice_export_timeout",
                details=details,
            ) from exc
        except OSError as exc:
            details = _libreoffice_details(
                cmd=lo_cmd,
                source_path=source,
                staged_input_path=staged_source,
                out_dir=out_dir,
                expected_pdf=expected_pdf,
                cwd=staging_dir,
                profile_dir=profile_dir,
                stderr=str(exc),
            )
            raise LibreOfficeExportError(
                "LibreOffice export failed before conversion completed",
                warning_code="libreoffice_export_failed",
                details=details,
            ) from exc

        if proc.returncode != 0:
            details = _libreoffice_details(
                cmd=lo_cmd,
                source_path=source,
                staged_input_path=staged_source,
                out_dir=out_dir,
                expected_pdf=expected_pdf,
                cwd=staging_dir,
                profile_dir=profile_dir,
                return_code=proc.returncode,
                stdout=proc.stdout,
                stderr=proc.stderr,
            )
            raise LibreOfficeExportError(
                "LibreOffice export returned a non-zero exit code",
                warning_code="libreoffice_export_return_code_nonzero",
                details=details,
            )

        selected_pdf = _select_libreoffice_pdf_output(source, out_dir)
        if selected_pdf and selected_pdf.exists() and selected_pdf.stat().st_size > 0:
            return selected_pdf

        warning_code = (
            "libreoffice_export_no_output_pdf"
            if not _pdf_output_candidates(out_dir)
            else "libreoffice_export_output_not_found"
        )
        details = _libreoffice_details(
            cmd=lo_cmd,
            source_path=source,
            staged_input_path=staged_source,
            out_dir=out_dir,
            expected_pdf=expected_pdf,
            cwd=staging_dir,
            profile_dir=profile_dir,
            return_code=proc.returncode,
            stdout=proc.stdout,
            stderr=proc.stderr,
        )
        raise LibreOfficeExportError(
            "LibreOffice export did not produce a usable PDF",
            warning_code=warning_code,
            details=details,
        )


def _write_stub(path: Path) -> None:
    """Write a 1×1 white PNG stub so the pipeline never gets a missing image."""
    if Image is not None:
        img = Image.new("RGB", (1920, 1080), color=(255, 255, 255))
        img.save(str(path), format="PNG")
    else:
        with open(path, "wb") as fh:
            fh.write(_WHITE_1X1_PNG)


def _make_text_image(text: str, out_path: Path, resolution: int = 1920) -> None:
    """Render *text* onto a white canvas and save as PNG."""
    if Image is None:
        _write_stub(out_path)
        return

    width = resolution
    height = int(resolution * 9 / 16)  # 16:9
    canvas = Image.new("RGB", (width, height), color=(255, 255, 255))
    draw = ImageDraw.Draw(canvas)

    font = None
    font_size = 32
    for fname in ("DejaVuSans.ttf", "arial.ttf", "Arial.ttf"):
        try:
            font = ImageFont.truetype(fname, font_size)
            break
        except Exception:
            pass
    if font is None:
        try:
            font = ImageFont.load_default()
        except Exception:
            font = None

    margin = 80
    max_w = width - 2 * margin
    line_h = font_size * 1.4 if font else 18

    # Wrap text
    lines: List[str] = []
    for paragraph in text.splitlines():
        if not paragraph.strip():
            lines.append("")
            continue
        words = paragraph.split()
        cur = ""
        for w in words:
            cand = f"{cur} {w}".strip()
            try:
                bbox = draw.textbbox((0, 0), cand, font=font) if font else (0, 0, len(cand) * 10, 20)
                w_px = bbox[2] - bbox[0]
            except Exception:
                w_px = len(cand) * 10
            if w_px <= max_w:
                cur = cand
            else:
                if cur:
                    lines.append(cur)
                cur = w
        if cur:
            lines.append(cur)

    y = margin
    for line in lines:
        if y + line_h > height - margin:
            break
        draw.text((margin, y), line, fill=(20, 20, 20), font=font)
        y += line_h

    canvas.save(str(out_path), format="PNG")


def _compact_text(text: str) -> str:
    """Normalize whitespace while preserving readable sentence flow."""
    return " ".join((text or "").replace("\xa0", " ").split())


# ---------------------------------------------------------------------------
# Strategy: LibreOffice → PDF → pdftoppm  (works for PPTX and DOCX)
# ---------------------------------------------------------------------------

def _export_via_libreoffice(source_path: str, out_dir: str, resolution: int = 1920) -> List[str]:
    out_dir_p = Path(out_dir)
    out_dir_p.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory() as tmpd:
        converted = _convert_via_libreoffice_to_pdf(source_path, Path(tmpd))

        # Convert resolution (px width) to approximate DPI (assume ~13in wide slide)
        dpi = max(72, resolution // 8)
        png_prefix = str(Path(tmpd) / "slide")
        _run(["pdftoppm", "-png", "-r", str(dpi), str(converted), png_prefix])

        produced = sorted(Path(tmpd).glob("slide-*.png")) or sorted(Path(tmpd).glob("*.png"))
        if not produced:
            raise RuntimeError("pdftoppm produced no PNGs")

        out_paths: List[str] = []
        for i, p in enumerate(produced, start=1):
            dst = out_dir_p / f"slide-{i}.png"
            _safe_replace(p, dst)
            out_paths.append(str(dst))
        return out_paths


# ---------------------------------------------------------------------------
# PPTX fallback: python-pptx + Pillow
# ---------------------------------------------------------------------------

def _wrap_text(text: str, draw, font, max_width: int) -> List[str]:
    if font is None:
        return text.splitlines()
    lines: List[str] = []
    for paragraph in text.splitlines():
        words = paragraph.split()
        if not words:
            lines.append("")
            continue
        cur = ""
        for w in words:
            cand = f"{cur} {w}".strip()
            try:
                bbox = draw.textbbox((0, 0), cand, font=font)
                width = bbox[2] - bbox[0]
            except Exception:
                # Older Pillow fallback
                width, _ = draw.textsize(cand, font=font)  # type: ignore[attr-defined]
            if width <= max_width:
                cur = cand
            else:
                if cur:
                    lines.append(cur)
                cur = w
        if cur:
            lines.append(cur)
    return lines


def _export_via_python_pptx(pptx_path: str, out_dir: str, resolution: int = 1920) -> List[str]:
    if Presentation is None or Image is None:
        raise RuntimeError("python-pptx and Pillow are required for PPTX fallback")

    prs = Presentation(pptx_path)
    slide_width = prs.slide_width or 9144000
    slide_height = prs.slide_height or 5143500
    scale = resolution / slide_width
    target_w = resolution
    target_h = int(slide_height * scale)

    out_dir_p = Path(out_dir)
    out_dir_p.mkdir(parents=True, exist_ok=True)

    font = None
    for fname in ("DejaVuSans.ttf", "arial.ttf", "Arial.ttf"):
        try:
            font = ImageFont.truetype(fname, 28)
            break
        except Exception:
            pass
    if font is None:
        try:
            font = ImageFont.load_default()
        except Exception:
            font = None

    exported: List[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        canvas = Image.new("RGB", (target_w, target_h), color=(255, 255, 255))
        draw = ImageDraw.Draw(canvas)

        for shape in slide.shapes:
            try:
                if getattr(shape, "shape_type", None) == getattr(MSO_SHAPE_TYPE, "PICTURE", None):
                    img = shape.image
                    blob = getattr(img, "blob", None)
                    if not blob:
                        continue
                    img_obj = Image.open(io.BytesIO(blob)).convert("RGBA")
                    left = int(getattr(shape, "left", 0) * scale)
                    top = int(getattr(shape, "top", 0) * scale)
                    w = max(1, int(getattr(shape, "width", img_obj.width) * scale))
                    h = max(1, int(getattr(shape, "height", img_obj.height) * scale))
                    img_obj = img_obj.resize((w, h), Image.LANCZOS)
                    canvas.paste(img_obj, (left, top), img_obj if img_obj.mode == "RGBA" else None)
            except Exception as e:
                logger.debug("Skipping picture on slide %d: %s", idx, e)

        for shape in slide.shapes:
            try:
                if hasattr(shape, "text_frame") and shape.text_frame is not None:
                    text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text)
                    if text.strip():
                        left = int(getattr(shape, "left", 0) * scale)
                        top = int(getattr(shape, "top", 0) * scale)
                        max_w = max(50, target_w - left - 20)
                        lines = _wrap_text(text, draw, font, max_w)
                        line_h = (getattr(font, "size", 14) if font else 14) * 1.2
                        for i, line in enumerate(lines):
                            draw.text((left + 5, top + int(i * line_h)), line, fill=(0, 0, 0), font=font)
                elif getattr(shape, "shape_type", None) == getattr(MSO_SHAPE_TYPE, "TABLE", None):
                    table = getattr(shape, "table", None)
                    if table:
                        left = int(getattr(shape, "left", 0) * scale)
                        top = int(getattr(shape, "top", 0) * scale)
                        max_w = max(50, target_w - left - 20)
                        y = top
                        for row in table.rows:
                            x = left
                            for cell in row.cells:
                                for i, line in enumerate(_wrap_text(cell.text or "", draw, font, max_w // max(1, len(row.cells)))):
                                    draw.text((x + 5, y + i * 14), line, fill=(0, 0, 0), font=font)
                                x += max_w // max(1, len(row.cells))
                            y += 24
            except Exception as e:
                logger.debug("Skipping text/table on slide %d: %s", idx, e)

        out_file = out_dir_p / f"slide-{idx}.png"
        try:
            canvas.save(out_file, format="PNG")
        except Exception:
            _write_stub(out_file)
        exported.append(str(out_file))

    return exported


def _shape_enum_value(name: str) -> Any:
    return getattr(MSO_SHAPE_TYPE, name, None) if MSO_SHAPE_TYPE is not None else None


def _shape_has_nonempty_text(shape: Any) -> bool:
    try:
        if not bool(getattr(shape, "has_text_frame", False)):
            return False
        text_frame = getattr(shape, "text_frame", None)
        if text_frame is None:
            return False
        return bool(str(getattr(text_frame, "text", "") or "").strip())
    except Exception:
        return False


def _shape_or_children_have_nonempty_text(shape: Any) -> bool:
    if _shape_has_nonempty_text(shape):
        return True
    try:
        return any(_shape_or_children_have_nonempty_text(child) for child in getattr(shape, "shapes", []) or [])
    except Exception:
        return False


def _shape_has_visible_fill_or_line(shape: Any) -> bool:
    try:
        fill = getattr(shape, "fill", None)
        fill_type = getattr(fill, "type", None)
        if fill_type is not None and str(fill_type).upper().find("BACKGROUND") < 0:
            return True
    except Exception:
        pass
    try:
        line = getattr(shape, "line", None)
        line_width = int(getattr(line, "width", 0) or 0)
        line_fill = getattr(line, "fill", None)
        if line_width > 0 or getattr(line_fill, "type", None) is not None:
            return True
    except Exception:
        pass
    return False


def _remove_shape_from_slide(shape: Any) -> bool:
    try:
        element = getattr(shape, "_element", None)
        parent = element.getparent() if element is not None else None
        if parent is None:
            return False
        parent.remove(element)
        return True
    except Exception:
        return False


def _clear_shape_text(shape: Any) -> bool:
    try:
        text_frame = getattr(shape, "text_frame", None)
        if text_frame is None:
            return False
        text_frame.clear()
        return True
    except Exception:
        try:
            shape.text = ""
            return True
        except Exception:
            return False


def _strip_text_from_pptx_copy(source_path: str, target_path: str) -> dict:
    """Best-effort PPTX text removal for source-background rendering."""
    if Presentation is None:
        raise RuntimeError("python-pptx is required for source background generation")

    target = Path(target_path)
    target.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(source_path, target)

    prs = Presentation(str(target))
    slide_warnings: list[list[str]] = [[] for _ in prs.slides]
    global_warnings: list[str] = []

    for slide_idx, slide in enumerate(prs.slides):
        shapes = list(getattr(slide, "shapes", []) or [])
        for shape in shapes:
            shape_type = getattr(shape, "shape_type", None)
            if shape_type == _shape_enum_value("GROUP") or hasattr(shape, "shapes"):
                if _shape_or_children_have_nonempty_text(shape):
                    slide_warnings[slide_idx].append("source_background_grouped_shape_skipped")
                continue
            if bool(getattr(shape, "has_table", False)) or shape_type == _shape_enum_value("TABLE"):
                if _shape_has_nonempty_text(shape) or bool(getattr(shape, "has_table", False)):
                    slide_warnings[slide_idx].append("source_background_table_text_skipped")
                continue
            if shape_type in {
                _shape_enum_value("CHART"),
                _shape_enum_value("DIAGRAM"),
                _shape_enum_value("IGX_GRAPHIC"),
            }:
                slide_warnings[slide_idx].append("source_background_chart_or_smartart_skipped")
                continue
            if not _shape_has_nonempty_text(shape):
                continue

            is_pure_text_box = (
                shape_type == _shape_enum_value("TEXT_BOX")
                and not _shape_has_visible_fill_or_line(shape)
            )
            if is_pure_text_box:
                if _remove_shape_from_slide(shape):
                    slide_warnings[slide_idx].append("source_background_text_removed")
                else:
                    slide_warnings[slide_idx].append("source_background_text_clear_failed")
                continue

            if _clear_shape_text(shape):
                slide_warnings[slide_idx].append("source_background_text_removed")
            else:
                slide_warnings[slide_idx].append("source_background_text_clear_failed")

    for warnings in slide_warnings:
        unique_slide_warnings = list(dict.fromkeys(warnings))
        warnings[:] = unique_slide_warnings
        global_warnings.extend(unique_slide_warnings)
        if any(warning.endswith("_skipped") for warning in unique_slide_warnings):
            warnings.append("source_background_partial_text_removal")
            global_warnings.append("source_background_partial_text_removal")

    prs.save(str(target))
    return {
        "pptx_path": str(target),
        "warnings": list(dict.fromkeys(global_warnings)),
        "slide_warnings": slide_warnings,
    }


def export_pptx_source_backgrounds(source_path: str, out_dir: str, resolution: int = 1920) -> dict:
    """Create cleaned PPTX slide backgrounds for the optional Source Background mode."""
    dependency_report = source_render_dependency_report()
    dependency_warnings = source_render_dependency_warnings(".pptx")
    if Path(source_path).suffix.lower() != ".pptx":
        return {
            "source_background_paths": [],
            "source_background_warnings": [],
            "source_background_slide_warnings": [],
            "source_background_details": [],
            "source_background_dependency_report": dependency_report,
        }

    _ensure_dir(out_dir)
    try:
        with tempfile.TemporaryDirectory(prefix="source-background-") as tmp:
            cleaned_path = Path(tmp) / Path(source_path).name
            strip_result = _strip_text_from_pptx_copy(source_path, str(cleaned_path))
            slide_warnings = list(strip_result.get("slide_warnings") or [])
            warnings = list(dict.fromkeys([*dependency_warnings, *list(strip_result.get("warnings") or [])]))
            details: list[dict[str, Any]] = []
            try:
                paths = _export_via_libreoffice(str(cleaned_path), out_dir, resolution)
                render_method = "libreoffice_pdf_raster"
            except Exception as exc:
                logger.warning(
                    "LibreOffice source background export failed for %s, using reconstructed fallback: %s",
                    source_path,
                    exc,
                )
                paths = _export_via_python_pptx(str(cleaned_path), out_dir, resolution)
                render_method = "python_pptx_reconstructed"
                details.extend(_libreoffice_failure_details(exc))
                warnings = list(
                    dict.fromkeys(
                        [
                            *warnings,
                            *_libreoffice_failure_warnings(exc),
                            "source_background_reconstructed",
                        ]
                    )
                )
            return {
                "source_background_paths": paths,
                "source_background_render_method": render_method,
                "source_background_warnings": warnings,
                "source_background_slide_warnings": slide_warnings,
                "source_background_details": details,
                "source_background_dependency_report": dependency_report,
            }
    except Exception as exc:
        logger.warning("PPTX source background generation failed for %s: %s", source_path, exc)
        return {
            "source_background_paths": [],
            "source_background_render_method": "failed",
            "source_background_warnings": list(dict.fromkeys([*dependency_warnings, "source_background_generation_failed"])),
            "source_background_slide_warnings": [],
            "source_background_details": [],
            "source_background_dependency_report": dependency_report,
        }


# ---------------------------------------------------------------------------
# PDF: pdftoppm or PyMuPDF
# ---------------------------------------------------------------------------

def _export_pdf_images(pdf_path: str, out_dir: str, resolution: int = 1920) -> List[str]:
    out_dir_p = Path(out_dir)
    out_dir_p.mkdir(parents=True, exist_ok=True)

    # Try pdftoppm
    try:
        png_prefix = str(out_dir_p / "slide")
        _run(["pdftoppm", "-png", pdf_path, png_prefix])
        produced = sorted(out_dir_p.glob("slide-*.png")) or sorted(out_dir_p.glob("*.png"))
        if produced:
            out_paths = []
            for i, p in enumerate(produced, start=1):
                dst = out_dir_p / f"slide-{i}.png"
                if p != dst:
                    p.replace(dst)
                out_paths.append(str(dst))
            return out_paths
    except Exception as e:
        logger.warning("pdftoppm failed for PDF: %s", e)

    # Try PyMuPDF
    if _HAVE_PYMUPDF:
        try:
            doc = _pymupdf.open(pdf_path)
            out_paths = []
            dpi = max(72, int(resolution / 11.7))  # approx for A4 at given width
            mat = _pymupdf.Matrix(dpi / 72, dpi / 72)
            for i, page in enumerate(doc, start=1):
                pix = page.get_pixmap(matrix=mat)
                dst = out_dir_p / f"slide-{i}.png"
                pix.save(str(dst))
                out_paths.append(str(dst))
            doc.close()
            return out_paths
        except Exception as e:
            logger.warning("PyMuPDF export failed: %s", e)

    raise RuntimeError(f"All PDF image export methods failed for {pdf_path}")


def _extract_pdf_text(pdf_path: str, notes_dir: Path) -> List[str]:
    """Extract text per page; return list of note file paths."""
    # Try PyMuPDF
    if _HAVE_PYMUPDF:
        try:
            doc = _pymupdf.open(pdf_path)
            out_paths = []
            for i, page in enumerate(doc, start=1):
                text = page.get_text().strip()
                if not text:
                    text = f"Page {i}."  # fallback if empty
                path = notes_dir / f"slide-{i}.txt"
                path.write_text(text, encoding="utf-8")
                out_paths.append(str(path))
            doc.close()
            return out_paths
        except Exception as e:
            logger.warning("PyMuPDF text extraction failed: %s", e)

    # Fallback: try pdftotext (from poppler-utils)
    try:
        with tempfile.TemporaryDirectory() as tmpd:
            txt_file = Path(tmpd) / "output.txt"
            _run(["pdftotext", "-layout", "-enc", "UTF-8", pdf_path, str(txt_file)])
            if txt_file.exists():
                full_text = txt_file.read_text(encoding="utf-8", errors="replace")
                # Split by form feed or approximate pages (pdftotext doesn't split by page easily)
                # For simplicity, if it's multi-page, we might need to use -f -l, but for now, assume single page or split manually
                # Actually, pdftotext without -layout might not preserve pages well.
                # Better: use pdftotext with -f 1 -l 1 for each page, but that's inefficient.
                # For now, if PyMuPDF failed, and pdftotext succeeds, use it as one block, but since we need per page, perhaps count pages first.
                # To keep it simple, use pdftotext and split by \f (form feed) if present.
                pages = full_text.split("\f") if "\f" in full_text else [full_text]
                out_paths = []
                for i, text in enumerate(pages, start=1):
                    text = text.strip() or f"Page {i}."
                    path = notes_dir / f"slide-{i}.txt"
                    path.write_text(text, encoding="utf-8")
                    out_paths.append(str(path))
                return out_paths
    except Exception as e:
        logger.warning("pdftotext fallback failed: %s", e)

    # Last resort: determine slide count and write stubs
    existing = sorted(notes_dir.parent.parent.glob("images/slide-*.png"))
    n = len(existing) or 1
    out_paths = []
    for i in range(1, n + 1):
        path = notes_dir / f"slide-{i}.txt"
        if not path.exists():
            path.write_text(f"Page {i}.", encoding="utf-8")
        out_paths.append(str(path))
    return out_paths


# ---------------------------------------------------------------------------
# DOCX: LibreOffice fallback, then python-docx + Pillow
# ---------------------------------------------------------------------------

def _docx_split_slides(doc) -> List[str]:
    """
    Split a python-docx document into slide/page-like text chunks.

    Strategy:
      - Keep heading paragraphs as structural boundaries.
      - Include normal paragraphs as body text.
      - Include table rows as pipe-separated lines so tabular content is not lost.
      - Split long runs into manageable chunks for TTS/subtitle pipeline.
    """
    slides: List[str] = []
    current: List[str] = []

    def _flush_current() -> None:
        if current:
            slides.append("\n".join(current).strip())
            current.clear()

    # Paragraphs (headings + body)
    for para in getattr(doc, "paragraphs", []):
        text = _compact_text(getattr(para, "text", ""))
        if not text:
            continue

        style_name = _compact_text(
            getattr(getattr(para, "style", None), "name", "")
        ).lower()
        is_heading = style_name.startswith("heading") or style_name in {"title", "subtitle"}

        if is_heading and current:
            _flush_current()

        current.append(f"# {text}" if is_heading else text)

        # Keep chunks readable and bounded for downstream TTS/subtitles
        if len(current) >= 12:
            _flush_current()

    # Tables (append as standalone chunks so table text remains available)
    for table in getattr(doc, "tables", []):
        table_lines: List[str] = []
        for row in getattr(table, "rows", []):
            cells: List[str] = []
            for cell in getattr(row, "cells", []):
                cell_text = _compact_text(getattr(cell, "text", ""))
                if cell_text:
                    cells.append(cell_text)
            if cells:
                table_lines.append(" | ".join(cells))
        if table_lines:
            _flush_current()
            slides.append("\n".join(table_lines))

    _flush_current()

    # If style parsing produced nothing, fall back to plain text splitting.
    if not slides:
        plain_lines = [
            _compact_text(getattr(p, "text", ""))
            for p in getattr(doc, "paragraphs", [])
        ]
        plain_lines = [line for line in plain_lines if line]
        if plain_lines:
            return _docx_split_slides_from_text("\n".join(plain_lines))

    return slides or ["(empty document)"]

def _export_docx_images_reconstructed(docx_path: str, out_dir: str, resolution: int = 1920) -> List[str]:
    # Attempt a richer Python-based rendering: prefer embedding any images
    # found in the DOCX; otherwise render the text onto a page-like canvas
    # (paper background + margin) instead of a bare whiteboard.
    if not _HAVE_DOCX or Image is None:
        raise RuntimeError("python-docx and Pillow required for DOCX image fallback")

    doc = _DocxDocument(docx_path)
    slides = _docx_split_slides(doc)
    out_dir_p = Path(out_dir)
    out_dir_p.mkdir(parents=True, exist_ok=True)
    out_paths: List[str] = []

    # Collect embedded images (blobs) from document parts (best-effort).
    embedded_images: List[bytes] = []
    try:
        for rel in getattr(doc.part, "_rels", {}).values():
            try:
                part = getattr(rel, "target_part", None) or getattr(rel, "_target", None)
                if part is None:
                    continue
                blob = getattr(part, "blob", None)
                if blob:
                    embedded_images.append(blob)
            except Exception:
                continue
    except Exception:
        embedded_images = []

    for i, text in enumerate(slides, start=1):
        dst = out_dir_p / f"slide-{i}.png"

        # If we have embedded images, paste the first one (scaled) onto a
        # page canvas and render the slide text below/over it. This produces
        # an image that more closely resembles a document page than a
        # pure whiteboard text render.
        try:
            if embedded_images:
                try:
                    img = Image.open(io.BytesIO(embedded_images[0])).convert("RGBA")
                    # If the embedded image is very low-res, fall back to text
                    if img.width < max(200, int(resolution * 0.4)):
                        raise ValueError("embedded image too low-res")

                    page_w = resolution
                    page_h = int(resolution * 11 / 8.5)  # taller for doc pages
                    canvas = Image.new("RGB", (page_w, page_h), color=(245, 245, 245))

                    # Paste image centered with margins
                    max_img_w = int(page_w * 0.9)
                    scale = min(1.0, max_img_w / img.width)
                    new_w = max(1, int(img.width * scale))
                    new_h = max(1, int(img.height * scale))
                    img_resized = img.resize((new_w, new_h), Image.LANCZOS)
                    left = (page_w - new_w) // 2
                    top = 60
                    canvas.paste(img_resized, (left, top), img_resized if img_resized.mode == "RGBA" else None)

                    # Render text area below image
                    draw = ImageDraw.Draw(canvas)
                    font = None
                    font_size = 20
                    for fname in ("DejaVuSans.ttf", "arial.ttf", "Arial.ttf"):
                        try:
                            font = ImageFont.truetype(fname, font_size)
                            break
                        except Exception:
                            pass
                    if font is None:
                        try:
                            font = ImageFont.load_default()
                        except Exception:
                            font = None

                    margin = 40
                    text_y = top + new_h + 20
                    max_w = page_w - 2 * margin
                    line_h = (getattr(font, "size", font_size) if font else font_size) * 1.3
                    # simple wrapping
                    lines = []
                    for paragraph in (text or f"Slide {i}").splitlines():
                        if not paragraph.strip():
                            lines.append("")
                            continue
                        words = paragraph.split()
                        cur = ""
                        for w in words:
                            cand = f"{cur} {w}".strip()
                            try:
                                bbox = draw.textbbox((0, 0), cand, font=font) if font else (0, 0, len(cand) * 8, 20)
                                w_px = bbox[2] - bbox[0]
                            except Exception:
                                w_px = len(cand) * 8
                            if w_px <= max_w:
                                cur = cand
                            else:
                                if cur:
                                    lines.append(cur)
                                cur = w
                        if cur:
                            lines.append(cur)

                    for line in lines:
                        if text_y + line_h > page_h - margin:
                            break
                        draw.text((margin, text_y), line, fill=(20, 20, 20), font=font)
                        text_y += line_h

                    canvas.save(str(dst), format="PNG")
                    out_paths.append(str(dst))
                    continue
                except Exception:
                    # Fall through to text rendering on error
                    logger.debug("Embedded-image rendering failed for %s, using text render", dst)

            # Default: render text onto a page-like canvas
            _make_text_image(text or f"Slide {i}", dst, resolution)
            out_paths.append(str(dst))
        except Exception:
            # Last resort: white stub
            _write_stub(dst)
            out_paths.append(str(dst))

    return out_paths or [str(out_dir_p / "slide-1.png")]


def _export_docx_images(docx_path: str, out_dir: str, resolution: int = 1920) -> List[str]:
    # LibreOffice handles DOCX the same way as PPTX.
    try:
        return _export_via_libreoffice(docx_path, out_dir, resolution)
    except Exception as e:
        logger.warning("LibreOffice failed for DOCX: %s", e)
    return _export_docx_images_reconstructed(docx_path, out_dir, resolution)


def _docx_split_slides_from_text(full_text: str) -> List[str]:
    """Split plain text into slide-sized chunks by headings or fixed count."""
    slides: List[str] = []
    current: List[str] = []
    for line in full_text.splitlines():
        line = line.strip()
        if not line:
            continue
        # Simple heuristic: lines starting with uppercase or numbers as headings
        is_heading = line[0].isupper() and len(line.split()) < 10  # rough
        if is_heading and current:
            slides.append("\n".join(current))
            current = []
        current.append(line)
        # Also split every ~10 lines
        if len(current) >= 10 and not is_heading:
            slides.append("\n".join(current))
            current = []
    if current:
        slides.append("\n".join(current))
    return slides or ["(empty document)"]


def _extract_docx_text(docx_path: str, notes_dir: Path) -> List[str]:
    # Prefer the same physical page model as the visual DOCX export:
    # DOCX -> PDF -> page text. This keeps text/page count aligned with
    # rasterized page images when LibreOffice is available.
    try:
        with tempfile.TemporaryDirectory() as tmpd:
            converted = _convert_via_libreoffice_to_pdf(docx_path, Path(tmpd))
            pdf_notes = _extract_pdf_text(str(converted), notes_dir)
            if pdf_notes:
                return pdf_notes
    except Exception as e:
        logger.warning("LibreOffice PDF text extraction failed for DOCX: %s", e)

    # Fallback to python-docx content extraction when page-based extraction is unavailable.
    if _HAVE_DOCX:
        try:
            doc = _DocxDocument(docx_path)
            slides = _docx_split_slides(doc)
            out_paths = []
            for i, text in enumerate(slides, start=1):
                path = notes_dir / f"slide-{i}.txt"
                path.write_text(text, encoding="utf-8")
                out_paths.append(str(path))
            return out_paths
        except Exception as e:
            logger.warning("python-docx text extraction failed: %s", e)

    # Fallback: try LibreOffice to convert to txt
    try:
        with (
            tempfile.TemporaryDirectory() as tmpd,
            tempfile.TemporaryDirectory(prefix="lo-profile-") as profile_tmp,
            tempfile.TemporaryDirectory(prefix="lo-input-") as stage_tmp,
        ):
            soffice = _find_soffice_executable() or "soffice"
            staged_docx = _stage_libreoffice_input(Path(docx_path).resolve(), Path(stage_tmp).resolve())
            _run([
                soffice, "--headless", "--nologo", "--nofirststartwizard", "--nolockcheck", "--nodefault",
                f"-env:UserInstallation={Path(profile_tmp).resolve().as_uri()}",
                "--convert-to", "txt:Text",
                "--outdir", tmpd,
                str(staged_docx),
            ])
            txt_candidates = sorted(Path(tmpd).glob("*.txt")) + sorted(Path(tmpd).glob("*.TXT"))
            if not txt_candidates:
                stem_txt = Path(tmpd) / f"{Path(docx_path).stem}.txt"
                if stem_txt.exists():
                    txt_candidates = [stem_txt]

            if txt_candidates:
                full_text = txt_candidates[0].read_text(encoding="utf-8", errors="replace")
                # Split into slides by headings or fixed count
                slides = _docx_split_slides_from_text(full_text)
                out_paths = []
                for i, text in enumerate(slides, start=1):
                    path = notes_dir / f"slide-{i}.txt"
                    path.write_text(text, encoding="utf-8")
                    out_paths.append(str(path))
                return out_paths
    except Exception as e:
        logger.warning("LibreOffice text extraction fallback failed: %s", e)

    # Last resort: create empty stubs matching image count
    existing = sorted(notes_dir.parent.parent.glob("images/slide-*.png"))
    n = len(existing) or 1
    out_paths = []
    for i in range(1, n + 1):
        path = notes_dir / f"slide-{i}.txt"
        path.write_text(f"Slide {i}.", encoding="utf-8")
        out_paths.append(str(path))
    return out_paths


# ---------------------------------------------------------------------------
# TXT: split by blank lines, render each block as image
# ---------------------------------------------------------------------------

def _txt_split_slides(txt_path: str) -> List[str]:
    """Split a plain-text file on double newlines into slide blocks."""
    raw = Path(txt_path).read_text(encoding="utf-8", errors="replace")
    blocks = [b.strip() for b in raw.split("\n\n") if b.strip()]
    return blocks or ["(empty file)"]


def _export_txt_images(txt_path: str, out_dir: str, resolution: int = 1920) -> List[str]:
    slides = _txt_split_slides(txt_path)
    out_dir_p = Path(out_dir)
    out_dir_p.mkdir(parents=True, exist_ok=True)
    out_paths = []
    for i, text in enumerate(slides, start=1):
        dst = out_dir_p / f"slide-{i}.png"
        _make_text_image(text, dst, resolution)
        out_paths.append(str(dst))
    return out_paths


def _extract_txt_text(txt_path: str, notes_dir: Path) -> List[str]:
    slides = _txt_split_slides(txt_path)
    out_paths = []
    for i, text in enumerate(slides, start=1):
        path = notes_dir / f"slide-{i}.txt"
        path.write_text(text, encoding="utf-8")
        out_paths.append(str(path))
    return out_paths


# ---------------------------------------------------------------------------
# Image sources: one uploaded image becomes one slide
# ---------------------------------------------------------------------------

def _flatten_image_to_rgb(image) -> Any:
    if Image is None:
        return image
    if ImageOps is not None:
        try:
            image = ImageOps.exif_transpose(image)
        except Exception:
            pass
    if image.mode in {"RGBA", "LA"} or (image.mode == "P" and "transparency" in getattr(image, "info", {})):
        rgba = image.convert("RGBA")
        canvas = Image.new("RGBA", rgba.size, (255, 255, 255, 255))
        canvas.alpha_composite(rgba)
        return canvas.convert("RGB")
    return image.convert("RGB")


def _export_image_source(source_path: str, out_dir: str) -> tuple[List[str], List[str]]:
    if Image is None:
        raise RuntimeError("Pillow is required for image source export")

    warnings: List[str] = []
    out_dir_p = Path(out_dir)
    out_dir_p.mkdir(parents=True, exist_ok=True)
    dst = out_dir_p / "slide-1.png"
    ext = Path(source_path).suffix.lower()
    with Image.open(source_path) as img:
        if ext == ".gif" and bool(getattr(img, "is_animated", False)):
            warnings.append("animated_gif_first_frame_only")
            try:
                img.seek(0)
            except Exception:
                pass
        frame = _flatten_image_to_rgb(img.copy())
        frame.save(dst, format="PNG")
    return [str(dst)], warnings


def _extract_image_text(image_path: str, notes_dir: Path) -> List[str]:
    notes_dir.mkdir(parents=True, exist_ok=True)
    out_path = notes_dir / "slide-1.txt"
    text = ""
    try:
        from worker.ai_agents.ocr_bridge import OCRBridge

        text = OCRBridge().extract_text(image_path=image_path) or ""
    except Exception as exc:
        logger.warning("Image OCR extraction skipped for %s: %s", image_path, exc)
    out_path.write_text(str(text or ""), encoding="utf-8")
    return [str(out_path)]


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def _fallback_stub_paths(source_path: str, out_dir: str, ext: str) -> List[str]:
    try:
        if ext == ".pptx" and Presentation is not None:
            n = len(Presentation(source_path).slides)
        else:
            n = 1
    except Exception:
        n = 1
    out_dir_p = Path(out_dir)
    out_dir_p.mkdir(parents=True, exist_ok=True)
    fallback_paths = []
    for i in range(1, n + 1):
        dst = out_dir_p / f"slide-{i}.png"
        _write_stub(dst)
        fallback_paths.append(str(dst))
    return fallback_paths


def _export_metadata_payload(
    image_paths: List[str],
    *,
    method: str,
    warnings: List[str] | None = None,
    details: list[dict[str, Any]] | None = None,
    dependency_report: dict | None = None,
) -> dict:
    return {
        "image_paths": image_paths,
        "source_render_method": method,
        "source_render_warnings": list(dict.fromkeys(warnings or [])),
        "source_render_details": list(details or []),
        "source_render_dependency_report": dict(dependency_report or source_render_dependency_report()),
    }


def export_slide_images_with_metadata(source_path: str, out_dir: str, resolution: int = 1920) -> dict:
    """Export slide/page images and include source-render fidelity metadata."""
    _ensure_dir(out_dir)
    ext = Path(source_path).suffix.lower()
    dependency_report = source_render_dependency_report()
    dependency_warnings = source_render_dependency_warnings(ext)

    try:
        if ext in _IMAGE_SOURCE_EXTS:
            paths, warnings = _export_image_source(source_path, out_dir)
            return _export_metadata_payload(
                paths,
                method="image_first_frame_png",
                warnings=warnings,
                dependency_report=dependency_report,
            )

        if ext == ".pdf":
            return _export_metadata_payload(
                _export_pdf_images(source_path, out_dir, resolution),
                method="pdf_raster",
                warnings=dependency_warnings,
                dependency_report=dependency_report,
            )

        if ext == ".docx":
            try:
                return _export_metadata_payload(
                    _export_via_libreoffice(source_path, out_dir, resolution),
                    method="libreoffice_pdf_raster",
                    warnings=dependency_warnings,
                    dependency_report=dependency_report,
                )
            except Exception as e:
                logger.warning("LibreOffice failed for DOCX, using reconstructed fallback: %s", e)
                return _export_metadata_payload(
                    _export_docx_images_reconstructed(source_path, out_dir, resolution),
                    method="python_docx_reconstructed",
                    warnings=[
                        *dependency_warnings,
                        *_libreoffice_failure_warnings(e),
                        "original_fidelity_reconstructed",
                    ],
                    details=_libreoffice_failure_details(e),
                    dependency_report=dependency_report,
                )

        if ext == ".txt":
            return _export_metadata_payload(
                _export_txt_images(source_path, out_dir, resolution),
                method="txt_whiteboard",
                dependency_report=dependency_report,
            )

        try:
            return _export_metadata_payload(
                _export_via_libreoffice(source_path, out_dir, resolution),
                method="libreoffice_pdf_raster",
                warnings=dependency_warnings,
                dependency_report=dependency_report,
            )
        except Exception as e:
            logger.warning("LibreOffice export failed, falling back to reconstructed python-pptx: %s", e)
            return _export_metadata_payload(
                _export_via_python_pptx(source_path, out_dir, resolution),
                method="python_pptx_reconstructed",
                warnings=[
                    *dependency_warnings,
                    *_libreoffice_failure_warnings(e),
                    "original_fidelity_reconstructed",
                ],
                details=_libreoffice_failure_details(e),
                dependency_report=dependency_report,
            )

    except Exception as exc:
        logger.exception("export_slide_images failed for %s: %s", source_path, exc)
        return _export_metadata_payload(
            _fallback_stub_paths(source_path, out_dir, ext),
            method="stub",
            warnings=[*dependency_warnings, "source_render_stub"],
            dependency_report=dependency_report,
        )


def export_slide_images(source_path: str, out_dir: str, resolution: int = 1920) -> List[str]:
    """Backward-compatible wrapper returning only exported image paths."""
    metadata = export_slide_images_with_metadata(source_path, out_dir, resolution)
    return list(metadata.get("image_paths") or [])


def extract_speaker_notes(source_path: str, out_dir: str) -> List[str]:
    """
    Extract narration text from *source_path* to *out_dir*/notes/.

    Dispatches on file extension (.pptx/.pdf/.docx/.txt).
    Returns list of note-file paths.
    """
    notes_dir = Path(out_dir) / "notes"
    notes_dir.mkdir(parents=True, exist_ok=True)
    ext = Path(source_path).suffix.lower()

    try:
        if ext in _IMAGE_SOURCE_EXTS:
            return _extract_image_text(source_path, notes_dir)

        if ext == ".pdf":
            return _extract_pdf_text(source_path, notes_dir)

        elif ext == ".docx":
            return _extract_docx_text(source_path, notes_dir)

        elif ext == ".txt":
            return _extract_txt_text(source_path, notes_dir)

        else:
            # PPTX
            if Presentation is None:
                raise RuntimeError("python-pptx is required to extract PPTX notes")
            prs = Presentation(source_path)
            out_paths: List[str] = []
            for i, slide in enumerate(prs.slides, start=1):
                notes_text = ""
                try:
                    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                        notes_text = slide.notes_slide.notes_text_frame.text or ""
                except Exception:
                    pass
                    
                if not notes_text.strip():
                    texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame") and shape.text_frame is not None:
                            t = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text)
                            if t.strip(): texts.append(t.strip())
                    notes_text = ". ".join(texts)
                    
                path = notes_dir / f"slide-{i}.txt"
                path.write_text(notes_text, encoding="utf-8")
                out_paths.append(str(path))
            return out_paths

    except Exception as exc:
        logger.exception("extract_speaker_notes failed for %s: %s", source_path, exc)
        # Return whatever note files already exist; create empties for missing ones
        existing_images = sorted(Path(out_dir).parent.glob("images/slide-*.png") if Path(out_dir).parent.exists() else [])
        n = len(existing_images) or 1
        out_paths = []
        for i in range(1, n + 1):
            path = notes_dir / f"slide-{i}.txt"
            if not path.exists():
                path.write_text(f"Slide {i}.", encoding="utf-8")
            out_paths.append(str(path))
        return out_paths


# ---------------------------------------------------------------------------
# CLI for manual testing
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    import argparse

    logging.basicConfig(level=logging.DEBUG, stream=sys.stdout)
    ap = argparse.ArgumentParser(description="Extract slides + notes from PPTX/PDF/DOCX/TXT")
    ap.add_argument("source", help="Path to .pptx/.pdf/.docx/.txt file")
    ap.add_argument("--images-dir", default="out/images", help="Output dir for PNGs")
    ap.add_argument("--notes-dir", default="out", help="Base dir for notes/ subdir")
    ap.add_argument("--resolution", type=int, default=1920)
    args = ap.parse_args()

    imgs = export_slide_images(args.source, args.images_dir, resolution=args.resolution)
    notes = extract_speaker_notes(args.source, args.notes_dir)
    print(f"Exported {len(imgs)} images, {len(notes)} note files.")
