# pyright: reportMissingImports=false

import os
import sys
from pathlib import Path

import django
import pytest

REPO_ROOT = Path(__file__).resolve().parents[2]
API_ROOT = REPO_ROOT / "services" / "api"
SERVICES_ROOT = REPO_ROOT / "services"
if str(API_ROOT) not in sys.path:
    sys.path.insert(0, str(API_ROOT))
if str(SERVICES_ROOT) not in sys.path:
    sys.path.insert(0, str(SERVICES_ROOT))

os.environ.setdefault("DJANGO_SETTINGS_MODULE", "config.settings")
django.setup()

from django.contrib.auth.models import User  # noqa: E402

from core.models import Project, TranscriptPage, UserProfile  # noqa: E402
from core.serializers import TranscriptPageSerializer  # noqa: E402
from scripts import pptx_extract  # noqa: E402
from worker import tasks as worker_tasks  # noqa: E402


PNG_1X1 = (
    b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
    b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc```\x00\x00"
    b"\x00\x04\x00\x01\xf6\x178U\x00\x00\x00\x00IEND\xaeB`\x82"
)


def _make_teacher(username: str) -> User:
    user = User.objects.create_user(username=username, password="pass")
    UserProfile.objects.create(user=user, role="teacher")
    return user


def _make_project(username: str) -> Project:
    return Project.objects.create(title=f"Extraction fidelity {username}", user=_make_teacher(username))


def _write_png(path: Path) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_bytes(PNG_1X1)
    return path


def _long_slide_text() -> str:
    return " ".join(
        f"Sentence {index} keeps the original source slide mapped to one visual frame."
        for index in range(80)
    )


def _sync_one_page(project: Project, tmp_path: Path, monkeypatch, *, source_type: str, page_key: str = "s1-p1") -> TranscriptPage:
    storage_root = tmp_path / "storage"
    image_path = _write_png(storage_root / str(project.id) / "images" / f"{source_type}-slide-1.png")
    monkeypatch.setattr(worker_tasks, "STORAGE_ROOT", str(storage_root))

    worker_tasks._sync_transcript_pages_from_export(
        project.id,
        [
            {
                "index": 0,
                "source_slide_index": 0,
                "split_index": 0,
                "page_key": page_key,
                "source_type": source_type,
                "image_path": str(image_path),
                "original_background_path": str(image_path),
                "original_text": "Visible page text",
                "narration_text": "Narrated page text",
                "subtitle_chunks": ["Narrated page text"],
                "whiteboard_mode": False,
            }
        ],
    )
    return TranscriptPage.objects.get(project=project, page_key=page_key)


def test_docx_text_extraction_prefers_pdf_page_text(tmp_path, monkeypatch):
    docx_path = tmp_path / "lesson.docx"
    docx_path.write_bytes(b"dummy docx")
    notes_dir = tmp_path / "notes"
    notes_dir.mkdir()

    def fake_convert(_source_path: str, out_dir: Path) -> Path:
        pdf_path = out_dir / "lesson.pdf"
        pdf_path.write_bytes(b"%PDF-1.4")
        return pdf_path

    def fake_extract_pdf_text(pdf_path: str, target_notes_dir: Path) -> list[str]:
        assert Path(pdf_path).name == "lesson.pdf"
        assert target_notes_dir == notes_dir
        out_paths = []
        for idx, text in enumerate(["Page one text", "Page two text"], start=1):
            path = target_notes_dir / f"slide-{idx}.txt"
            path.write_text(text, encoding="utf-8")
            out_paths.append(str(path))
        return out_paths

    monkeypatch.setattr(pptx_extract, "_convert_via_libreoffice_to_pdf", fake_convert)
    monkeypatch.setattr(pptx_extract, "_extract_pdf_text", fake_extract_pdf_text)

    paths = pptx_extract._extract_docx_text(str(docx_path), notes_dir)

    assert len(paths) == 2
    assert [Path(path).read_text(encoding="utf-8") for path in paths] == ["Page one text", "Page two text"]


@pytest.mark.django_db
def test_docx_export_uses_raster_page_count_and_handles_text_mismatch(tmp_path, monkeypatch):
    project = _make_project("docx_page_count")
    storage_root = tmp_path / "storage"
    source_path = tmp_path / "source.docx"
    source_path.write_bytes(b"fake docx")
    monkeypatch.setattr(worker_tasks, "STORAGE_ROOT", str(storage_root))

    def fake_export_slide_images_with_metadata(_source_path: str, out_dir: str) -> dict:
        image_paths = [
            str(_write_png(Path(out_dir) / "slide-1.png")),
            str(_write_png(Path(out_dir) / "slide-2.png")),
            str(_write_png(Path(out_dir) / "slide-3.png")),
        ]
        return {
            "image_paths": image_paths,
            "source_render_method": "libreoffice_pdf_raster",
            "source_render_warnings": [],
        }

    def fake_extract_speaker_notes(_source_path: str, out_dir: str) -> list[str]:
        note_path = Path(out_dir) / "notes" / "slide-1.txt"
        note_path.parent.mkdir(parents=True, exist_ok=True)
        note_path.write_text("Only the first physical page has extracted text.", encoding="utf-8")
        return [str(note_path)]

    monkeypatch.setattr(pptx_extract, "export_slide_images_with_metadata", fake_export_slide_images_with_metadata)
    monkeypatch.setattr(pptx_extract, "extract_speaker_notes", fake_extract_speaker_notes)
    monkeypatch.setattr(worker_tasks.export_project, "update_state", lambda *args, **kwargs: None)

    slides = worker_tasks.export_project.run(str(project.id), str(source_path), False)

    assert len(slides) == 3
    assert [slide["source_type"] for slide in slides] == ["docx", "docx", "docx"]
    assert slides[0]["original_text"] == "Only the first physical page has extracted text."
    assert slides[1]["original_text"] == ""
    assert slides[2]["narration_text"] == ""
    assert all(slide["original_background_path"] == slide["image_path"] for slide in slides)


@pytest.mark.django_db
def test_pptx_source_background_export_preserves_source_slide_mapping_for_long_text(tmp_path, monkeypatch):
    project = _make_project("pptx_source_mapping")
    storage_root = tmp_path / "storage"
    source_path = tmp_path / "source.pptx"
    source_path.write_bytes(b"fake pptx")
    note_text = _long_slide_text()
    monkeypatch.setattr(worker_tasks, "STORAGE_ROOT", str(storage_root))
    monkeypatch.setattr(worker_tasks.export_project, "update_state", lambda *args, **kwargs: None)

    def fake_export_slide_images_with_metadata(_source_path: str, out_dir: str) -> dict:
        return {
            "image_paths": [str(_write_png(Path(out_dir) / "slide-1.png"))],
            "source_render_method": "libreoffice_pdf_raster",
            "source_render_warnings": [],
        }

    def fake_source_backgrounds(_source_path: str, out_dir: str) -> dict:
        return {
            "source_background_paths": [str(_write_png(Path(out_dir) / "slide-1.png"))],
            "source_background_warnings": ["source_background_text_removed"],
            "source_background_slide_warnings": [],
        }

    def fake_extract_speaker_notes(_source_path: str, out_dir: str) -> list[str]:
        note_path = Path(out_dir) / "notes" / "slide-1.txt"
        note_path.parent.mkdir(parents=True, exist_ok=True)
        note_path.write_text(note_text, encoding="utf-8")
        return [str(note_path)]

    monkeypatch.setattr(pptx_extract, "export_slide_images_with_metadata", fake_export_slide_images_with_metadata)
    monkeypatch.setattr(pptx_extract, "export_pptx_source_backgrounds", fake_source_backgrounds)
    monkeypatch.setattr(pptx_extract, "extract_speaker_notes", fake_extract_speaker_notes)

    slides = worker_tasks.export_project.run(str(project.id), str(source_path), False)

    assert len(slides) == 1
    assert slides[0]["source_slide_index"] == 0
    assert slides[0]["split_index"] == 0
    assert slides[0]["page_key"] == "s1-p1"
    assert slides[0]["source_background_path"].replace("\\", "/").endswith("source_backgrounds/slide-1.png")
    assert slides[0]["original_text"] == note_text
    assert len(slides[0]["subtitle_chunks"]) > 1


@pytest.mark.django_db
def test_whiteboard_export_keeps_long_text_visual_splitting(tmp_path, monkeypatch):
    project = _make_project("whiteboard_source_splitting")
    storage_root = tmp_path / "storage"
    source_path = tmp_path / "source.pptx"
    source_path.write_bytes(b"fake pptx")
    note_text = _long_slide_text()
    monkeypatch.setattr(worker_tasks, "STORAGE_ROOT", str(storage_root))
    monkeypatch.setattr(worker_tasks.export_project, "update_state", lambda *args, **kwargs: None)

    def fake_export_slide_images_with_metadata(_source_path: str, out_dir: str) -> dict:
        return {
            "image_paths": [str(_write_png(Path(out_dir) / "slide-1.png"))],
            "source_render_method": "libreoffice_pdf_raster",
            "source_render_warnings": [],
        }

    def fake_extract_speaker_notes(_source_path: str, out_dir: str) -> list[str]:
        note_path = Path(out_dir) / "notes" / "slide-1.txt"
        note_path.parent.mkdir(parents=True, exist_ok=True)
        note_path.write_text(note_text, encoding="utf-8")
        return [str(note_path)]

    monkeypatch.setattr(pptx_extract, "export_slide_images_with_metadata", fake_export_slide_images_with_metadata)
    monkeypatch.setattr(pptx_extract, "export_pptx_source_backgrounds", lambda *_args, **_kwargs: {})
    monkeypatch.setattr(pptx_extract, "extract_speaker_notes", fake_extract_speaker_notes)

    slides = worker_tasks.export_project.run(str(project.id), str(source_path), True)

    assert len(slides) > 1
    assert {slide["source_slide_index"] for slide in slides} == {0}
    assert [slide["split_index"] for slide in slides] == list(range(len(slides)))
    assert all(slide["whiteboard_mode"] is True for slide in slides)


@pytest.mark.django_db
@pytest.mark.parametrize("source_type", ["docx", "pdf", "pptx"])
def test_visual_source_pages_default_to_original_background(tmp_path, monkeypatch, source_type):
    project = _make_project(f"{source_type}_original_background")

    page = _sync_one_page(project, tmp_path, monkeypatch, source_type=source_type)

    scene = page.editor_document["scene"]
    assert scene["background_mode"] == "original"
    assert scene["original_background_path"].endswith(f"{source_type}-slide-1.png")
    serialized_scene = TranscriptPageSerializer(page).data["editor_document"]["scene"]
    assert serialized_scene["has_original_background"] is True
    assert serialized_scene["original_background_url"]
    assert "original_background_path" not in serialized_scene


@pytest.mark.django_db
def test_transcript_resync_preserves_existing_custom_background(tmp_path, monkeypatch):
    project = _make_project("preserve_custom")
    page = TranscriptPage.objects.create(
        project=project,
        order=0,
        source_slide_index=0,
        split_index=0,
        page_key="s1-p1",
        original_text="Existing text",
        narration_text="Existing narration",
        editor_document={
            "version": 1,
            "scene": {
                "background_mode": "custom",
                "custom_background_path": f"uploads/{project.id}/backgrounds/custom.png",
                "background_fit": "cover",
                "text_scale": 1.25,
            },
        },
    )
    storage_root = tmp_path / "storage"
    image_path = _write_png(storage_root / str(project.id) / "images" / "docx-slide-1.png")
    monkeypatch.setattr(worker_tasks, "STORAGE_ROOT", str(storage_root))

    worker_tasks._sync_transcript_pages_from_export(
        project.id,
        [{"index": 0, "page_key": page.page_key, "source_type": "docx", "image_path": str(image_path)}],
    )

    page.refresh_from_db()
    scene = page.editor_document["scene"]
    assert scene["background_mode"] == "custom"
    assert scene["custom_background_path"] == f"uploads/{project.id}/backgrounds/custom.png"
    assert scene["background_fit"] == "cover"
    assert scene["text_scale"] == 1.25
    assert scene["original_background_path"] == f"{project.id}/images/docx-slide-1.png"


@pytest.mark.django_db
def test_transcript_resync_preserves_existing_whiteboard_mode(tmp_path, monkeypatch):
    project = _make_project("preserve_whiteboard")
    page = TranscriptPage.objects.create(
        project=project,
        order=0,
        source_slide_index=0,
        split_index=0,
        page_key="s1-p1",
        original_text="Existing text",
        narration_text="Existing narration",
        whiteboard_mode=True,
        editor_document={"version": 1, "scene": {"background_mode": "whiteboard"}},
    )
    storage_root = tmp_path / "storage"
    image_path = _write_png(storage_root / str(project.id) / "images" / "pdf-slide-1.png")
    monkeypatch.setattr(worker_tasks, "STORAGE_ROOT", str(storage_root))

    worker_tasks._sync_transcript_pages_from_export(
        project.id,
        [{"index": 0, "page_key": page.page_key, "source_type": "pdf", "image_path": str(image_path)}],
    )

    page.refresh_from_db()
    assert page.whiteboard_mode is True
    assert page.editor_document["scene"]["background_mode"] == "whiteboard"
    assert page.editor_document["scene"]["original_background_path"] == f"{project.id}/images/pdf-slide-1.png"


@pytest.mark.django_db
def test_pptx_sync_stores_source_background_without_defaulting_to_it(tmp_path, monkeypatch):
    project = _make_project("pptx_source_background")
    storage_root = tmp_path / "storage"
    original_path = _write_png(storage_root / str(project.id) / "images" / "slide-1.png")
    source_background_path = _write_png(storage_root / str(project.id) / "source_backgrounds" / "slide-1.png")
    monkeypatch.setattr(worker_tasks, "STORAGE_ROOT", str(storage_root))

    worker_tasks._sync_transcript_pages_from_export(
        project.id,
        [
            {
                "index": 0,
                "source_slide_index": 0,
                "page_key": "s1-p1",
                "source_type": "pptx",
                "image_path": str(original_path),
                "original_background_path": str(original_path),
                "source_background_path": str(source_background_path),
                "source_background_warnings": ["source_background_text_removed"],
                "original_text": "Extracted text remains editable",
                "narration_text": "Extracted text remains editable",
                "whiteboard_mode": False,
            }
        ],
    )

    page = TranscriptPage.objects.get(project=project, page_key="s1-p1")
    scene = page.editor_document["scene"]
    assert scene["background_mode"] == "original"
    assert scene["source_background_path"] == f"{project.id}/source_backgrounds/slide-1.png"
    assert scene["source_background_warnings"] == ["source_background_text_removed"]
    serialized_scene = TranscriptPageSerializer(page).data["editor_document"]["scene"]
    assert serialized_scene["background_mode"] == "original"
    assert serialized_scene["has_source_background"] is True
    assert serialized_scene["source_background_url"]
    assert "source_background_path" not in serialized_scene


@pytest.mark.django_db
def test_txt_source_defaults_to_whiteboard_without_original_background(tmp_path, monkeypatch):
    project = _make_project("txt_whiteboard")
    storage_root = tmp_path / "storage"
    image_path = _write_png(storage_root / str(project.id) / "images" / "slide-1.png")
    monkeypatch.setattr(worker_tasks, "STORAGE_ROOT", str(storage_root))

    worker_tasks._sync_transcript_pages_from_export(
        project.id,
        [
            {
                "index": 0,
                "source_slide_index": 0,
                "page_key": "s1-p1",
                "source_type": "txt",
                "image_path": str(image_path),
                "original_text": "Plain text source",
                "narration_text": "Plain text source",
                "whiteboard_mode": True,
            }
        ],
    )

    page = TranscriptPage.objects.get(project=project, page_key="s1-p1")
    scene = page.editor_document["scene"]
    assert scene["background_mode"] == "whiteboard"
    assert "original_background_path" not in scene
    serialized_scene = TranscriptPageSerializer(page).data["editor_document"]["scene"]
    assert serialized_scene["has_original_background"] is False
    assert serialized_scene["original_background_url"] == ""
    assert "original_background_path" not in serialized_scene


def test_reconstructed_office_export_records_fidelity_warning(tmp_path, monkeypatch):
    source_path = tmp_path / "lesson.pptx"
    source_path.write_bytes(b"fake pptx")

    def fail_libreoffice(*_args, **_kwargs):
        raise pptx_extract.LibreOfficeExportError(
            "libreoffice unavailable",
            warning_code="libreoffice_export_no_output_pdf",
            details={
                "warning_code": "libreoffice_export_no_output_pdf",
                "actual_output_files": ["conversion.log (12 bytes)"],
            },
        )

    def fake_python_pptx(_source_path: str, out_dir: str, _resolution: int = 1920) -> list[str]:
        return [str(_write_png(Path(out_dir) / "slide-1.png"))]

    monkeypatch.setattr(pptx_extract, "_export_via_libreoffice", fail_libreoffice)
    monkeypatch.setattr(pptx_extract, "_export_via_python_pptx", fake_python_pptx)

    metadata = pptx_extract.export_slide_images_with_metadata(str(source_path), str(tmp_path / "images"))

    assert metadata["source_render_method"] == "python_pptx_reconstructed"
    assert "libreoffice_export_failed" in metadata["source_render_warnings"]
    assert "libreoffice_export_no_output_pdf" in metadata["source_render_warnings"]
    assert "original_fidelity_reconstructed" in metadata["source_render_warnings"]
    assert metadata["source_render_details"][0]["actual_output_files"] == ["conversion.log (12 bytes)"]
    assert metadata["image_paths"][0].endswith("slide-1.png")


def test_docx_reconstructed_export_records_fidelity_warning(tmp_path, monkeypatch):
    source_path = tmp_path / "lesson.docx"
    source_path.write_bytes(b"fake docx")

    def fail_libreoffice(*_args, **_kwargs):
        raise RuntimeError("libreoffice unavailable")

    def fake_reconstructed(_source_path: str, out_dir: str, _resolution: int = 1920) -> list[str]:
        return [str(_write_png(Path(out_dir) / "slide-1.png"))]

    monkeypatch.setattr(pptx_extract, "_export_via_libreoffice", fail_libreoffice)
    monkeypatch.setattr(pptx_extract, "_export_docx_images_reconstructed", fake_reconstructed)

    metadata = pptx_extract.export_slide_images_with_metadata(str(source_path), str(tmp_path / "images"))

    assert metadata["source_render_method"] == "python_docx_reconstructed"
    assert "original_fidelity_reconstructed" in metadata["source_render_warnings"]
    assert metadata["image_paths"][0].endswith("slide-1.png")


def test_image_source_exports_single_original_slide_and_keeps_ocr_text(tmp_path, monkeypatch):
    source_path = _write_png(tmp_path / "lesson.png")
    storage_root = tmp_path / "storage"
    monkeypatch.setattr(worker_tasks, "STORAGE_ROOT", str(storage_root))
    monkeypatch.setattr(worker_tasks.export_project, "update_state", lambda *args, **kwargs: None)

    from worker.ai_agents import ocr_bridge

    monkeypatch.setattr(ocr_bridge.OCRBridge, "extract_text", lambda *_args, **_kwargs: "Image OCR narration")

    slides = worker_tasks.export_project.run("77", str(source_path), False)

    assert len(slides) == 1
    assert slides[0]["source_type"] == "png"
    assert slides[0]["source_render_method"] == "image_first_frame_png"
    assert slides[0]["source_render_warnings"] == []
    assert slides[0]["whiteboard_mode"] is False
    assert slides[0]["original_background_path"] == slides[0]["image_path"]
    assert slides[0]["original_text"] == "Image OCR narration"
    assert slides[0]["narration_text"] == "Image OCR narration"

    visual_assets = worker_tasks._visual_slide_assets_from_export(slides)
    assert visual_assets[0]["image_path"] == slides[0]["image_path"]


def test_animated_gif_source_normalizes_first_frame_with_warning(tmp_path):
    if pptx_extract.Image is None:
        pytest.skip("Pillow is not available")

    source_path = tmp_path / "lesson.gif"
    first = pptx_extract.Image.new("RGB", (4, 4), color=(255, 0, 0))
    second = pptx_extract.Image.new("RGB", (4, 4), color=(0, 0, 255))
    first.save(source_path, format="GIF", save_all=True, append_images=[second], duration=50, loop=0)

    metadata = pptx_extract.export_slide_images_with_metadata(str(source_path), str(tmp_path / "images"))

    assert metadata["source_render_method"] == "image_first_frame_png"
    assert metadata["source_render_warnings"] == ["animated_gif_first_frame_only"]
    assert len(metadata["image_paths"]) == 1
    assert metadata["image_paths"][0].endswith("slide-1.png")


def test_pptx_source_background_text_hiding_is_best_effort(tmp_path):
    if pptx_extract.Presentation is None:
        pytest.skip("python-pptx is not available")

    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    source_path = tmp_path / "source.pptx"
    cleaned_path = tmp_path / "cleaned.pptx"
    prs = pptx_extract.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(0.5))
    text_box.text = "Pure text box"
    filled_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.25), Inches(2), Inches(0.75))
    filled_shape.text = "Text in visible shape"
    filled_shape.fill.solid()
    filled_shape.fill.fore_color.rgb = RGBColor(230, 240, 255)
    table_shape = slide.shapes.add_table(1, 1, Inches(0.5), Inches(2.25), Inches(2), Inches(0.6))
    table_shape.table.cell(0, 0).text = "Table text"
    prs.save(source_path)

    result = pptx_extract._strip_text_from_pptx_copy(str(source_path), str(cleaned_path))
    cleaned = pptx_extract.Presentation(str(cleaned_path))
    cleaned_text = "\n".join(
        str(getattr(shape, "text", "") or "")
        for shape in cleaned.slides[0].shapes
        if bool(getattr(shape, "has_text_frame", False))
    )

    assert "Pure text box" not in cleaned_text
    assert "Text in visible shape" not in cleaned_text
    assert "source_background_text_removed" in result["warnings"]
    assert "source_background_table_text_skipped" in result["warnings"]
    assert "source_background_partial_text_removal" in result["warnings"]


def test_pptx_source_background_text_hiding_preserves_pictures_and_original(tmp_path):
    if pptx_extract.Presentation is None or pptx_extract.Image is None:
        pytest.skip("python-pptx and Pillow are not available")

    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Inches

    picture_path = tmp_path / "picture.png"
    pptx_extract.Image.new("RGB", (24, 24), color=(20, 90, 160)).save(picture_path)
    source_path = tmp_path / "source-with-picture.pptx"
    cleaned_path = tmp_path / "cleaned-with-picture.pptx"

    prs = pptx_extract.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(picture_path), Inches(0.25), Inches(0.25), Inches(1.5), Inches(1.5))
    text_box = slide.shapes.add_textbox(Inches(2), Inches(0.5), Inches(2), Inches(0.75))
    text_box.text = "Remove this text"
    prs.save(source_path)

    result = pptx_extract._strip_text_from_pptx_copy(str(source_path), str(cleaned_path))

    original = pptx_extract.Presentation(str(source_path))
    cleaned = pptx_extract.Presentation(str(cleaned_path))
    original_shapes = list(original.slides[0].shapes)
    cleaned_shapes = list(cleaned.slides[0].shapes)
    original_text = "\n".join(str(getattr(shape, "text", "") or "") for shape in original_shapes)
    cleaned_text = "\n".join(str(getattr(shape, "text", "") or "") for shape in cleaned_shapes)
    original_picture_count = sum(
        1 for shape in original_shapes if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE
    )
    cleaned_picture_count = sum(
        1 for shape in cleaned_shapes if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE
    )

    assert original_picture_count == 1
    assert cleaned_picture_count == 1
    assert "Remove this text" in original_text
    assert "Remove this text" not in cleaned_text
    assert "source_background_text_removed" in result["warnings"]


def test_source_render_dependency_warnings_surface_missing_tools(monkeypatch):
    for env_name in pptx_extract._LIBREOFFICE_CONFIG_ENV_VARS:
        monkeypatch.delenv(env_name, raising=False)
    monkeypatch.delenv("PROGRAMFILES", raising=False)
    monkeypatch.delenv("PROGRAMFILES(X86)", raising=False)
    monkeypatch.setattr(pptx_extract.shutil, "which", lambda _name: None)

    warnings = pptx_extract.source_render_dependency_warnings(".pptx")
    report = pptx_extract.source_render_dependency_report()

    assert "slide_render_dependency_missing_libreoffice" in warnings
    assert "slide_render_dependency_missing_pdftoppm" in warnings
    assert report["libreoffice_available"] is False
    assert report["pdftoppm_available"] is False
    assert "pymupdf_available" in report
    assert "python_pptx_available" in report


def test_libreoffice_nonzero_exit_records_diagnostic_warning_details(tmp_path, monkeypatch):
    source_path = tmp_path / "lesson.pptx"
    source_path.write_bytes(b"fake pptx")
    out_dir = tmp_path / "lo-out"
    out_dir.mkdir()
    (out_dir / "conversion.log").write_text("debug", encoding="utf-8")

    def fake_run(cmd, **_kwargs):
        return pptx_extract.subprocess.CompletedProcess(
            cmd,
            23,
            stdout="stdout-prefix-" + ("o" * 1400),
            stderr="stderr-prefix-" + ("e" * 1400),
        )

    monkeypatch.setattr(
        pptx_extract,
        "_find_soffice_executable",
        lambda: r"C:\Program Files\LibreOffice\program\soffice.exe",
    )
    monkeypatch.setattr(pptx_extract.subprocess, "run", fake_run)

    with pytest.raises(pptx_extract.LibreOfficeExportError) as raised:
        pptx_extract._convert_via_libreoffice_to_pdf(str(source_path), out_dir)

    exc = raised.value
    assert exc.warnings == ["libreoffice_export_failed", "libreoffice_export_return_code_nonzero"]
    details = exc.details
    assert details["warning_code"] == "libreoffice_export_return_code_nonzero"
    assert details["return_code"] == 23
    assert "soffice.exe" in details["command"]
    assert "--headless" in details["command"]
    assert "-env:UserInstallation=file:" in details["command"]
    assert details["stdout_tail"].endswith("o" * 1200)
    assert details["stderr_tail"].endswith("e" * 1200)
    assert details["expected_output_path"].endswith("lesson.pdf")
    assert any("conversion.log" in item for item in details["actual_output_files"])
    assert details["input_path"] == str(source_path.resolve())
    assert details["output_directory"] == str(out_dir.resolve())
    assert details["working_directory"] == str(tmp_path.resolve())


def test_libreoffice_missing_output_pdf_records_outdir_files(tmp_path, monkeypatch):
    source_path = tmp_path / "lesson.pptx"
    source_path.write_bytes(b"fake pptx")
    out_dir = tmp_path / "lo-out"
    out_dir.mkdir()
    (out_dir / "notes.txt").write_text("not a pdf", encoding="utf-8")

    def fake_run(cmd, **_kwargs):
        return pptx_extract.subprocess.CompletedProcess(cmd, 0, stdout="converted", stderr="")

    monkeypatch.setattr(pptx_extract, "_find_soffice_executable", lambda: "soffice")
    monkeypatch.setattr(pptx_extract.subprocess, "run", fake_run)

    with pytest.raises(pptx_extract.LibreOfficeExportError) as raised:
        pptx_extract._convert_via_libreoffice_to_pdf(str(source_path), out_dir)

    exc = raised.value
    assert "libreoffice_export_no_output_pdf" in exc.warnings
    assert exc.details["warning_code"] == "libreoffice_export_no_output_pdf"
    assert exc.details["return_code"] == 0
    assert any("notes.txt" in item for item in exc.details["actual_output_files"])


def test_libreoffice_empty_output_pdf_records_output_not_found(tmp_path, monkeypatch):
    source_path = tmp_path / "lesson.pptx"
    source_path.write_bytes(b"fake pptx")
    out_dir = tmp_path / "lo-out"
    out_dir.mkdir()

    def fake_run(cmd, **_kwargs):
        (out_dir / "lesson.pdf").write_bytes(b"")
        return pptx_extract.subprocess.CompletedProcess(cmd, 0, stdout="", stderr="")

    monkeypatch.setattr(pptx_extract, "_find_soffice_executable", lambda: "soffice")
    monkeypatch.setattr(pptx_extract.subprocess, "run", fake_run)

    with pytest.raises(pptx_extract.LibreOfficeExportError) as raised:
        pptx_extract._convert_via_libreoffice_to_pdf(str(source_path), out_dir)

    exc = raised.value
    assert "libreoffice_export_output_not_found" in exc.warnings
    assert exc.details["warning_code"] == "libreoffice_export_output_not_found"
    assert any("lesson.pdf (0 bytes)" in item for item in exc.details["actual_output_files"])


def test_pptx_source_background_uses_reconstructed_fallback_when_libreoffice_fails(tmp_path, monkeypatch):
    if pptx_extract.Presentation is None:
        pytest.skip("python-pptx is not available")

    source_path = tmp_path / "lesson.pptx"
    out_dir = tmp_path / "source_backgrounds"

    prs = pptx_extract.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    text_box = slide.shapes.add_textbox(0, 0, 1000000, 1000000)
    text_box.text = "Remove this"
    prs.save(source_path)

    def fail_libreoffice(*_args, **_kwargs):
        raise pptx_extract.LibreOfficeExportError(
            "libreoffice missing",
            warning_code="libreoffice_export_no_output_pdf",
            details={
                "warning_code": "libreoffice_export_no_output_pdf",
                "actual_output_files": ["notes.txt (9 bytes)"],
            },
        )

    def fake_python_pptx(_source_path: str, output_dir: str, _resolution: int = 1920) -> list[str]:
        return [str(_write_png(Path(output_dir) / "slide-1.png"))]

    monkeypatch.setattr(pptx_extract, "_export_via_libreoffice", fail_libreoffice)
    monkeypatch.setattr(pptx_extract, "_export_via_python_pptx", fake_python_pptx)

    metadata = pptx_extract.export_pptx_source_backgrounds(str(source_path), str(out_dir))

    assert len(metadata["source_background_paths"]) == 1
    assert metadata["source_background_render_method"] == "python_pptx_reconstructed"
    assert "libreoffice_export_failed" in metadata["source_background_warnings"]
    assert "libreoffice_export_no_output_pdf" in metadata["source_background_warnings"]
    assert "source_background_reconstructed" in metadata["source_background_warnings"]
    assert metadata["source_background_details"][0]["actual_output_files"] == ["notes.txt (9 bytes)"]
    assert "source_background_generation_failed" not in metadata["source_background_warnings"]


def test_pptx_source_background_generation_failure_records_warning(tmp_path, monkeypatch):
    source_path = tmp_path / "lesson.pptx"
    source_path.write_bytes(b"fake pptx")

    def fail_strip(*_args, **_kwargs):
        raise RuntimeError("unsupported deck")

    monkeypatch.setattr(pptx_extract, "_strip_text_from_pptx_copy", fail_strip)

    metadata = pptx_extract.export_pptx_source_backgrounds(str(source_path), str(tmp_path / "source_backgrounds"))

    assert metadata["source_background_paths"] == []
    assert "source_background_generation_failed" in metadata["source_background_warnings"]


@pytest.mark.django_db
def test_export_project_includes_pptx_source_background_metadata(tmp_path, monkeypatch):
    project = _make_project("pptx_source_background_export")
    storage_root = tmp_path / "storage"
    source_path = tmp_path / "lesson.pptx"
    source_path.write_bytes(b"fake pptx")
    monkeypatch.setattr(worker_tasks, "STORAGE_ROOT", str(storage_root))
    monkeypatch.setattr(worker_tasks.export_project, "update_state", lambda *args, **kwargs: None)

    def fake_export_slide_images_with_metadata(_source_path: str, out_dir: str) -> dict:
        return {
            "image_paths": [str(_write_png(Path(out_dir) / "slide-1.png"))],
            "source_render_method": "libreoffice_pdf_raster",
            "source_render_warnings": [],
            "source_render_details": [{"warning_code": "original_detail"}],
        }

    def fake_source_backgrounds(_source_path: str, out_dir: str) -> dict:
        return {
            "source_background_paths": [str(_write_png(Path(out_dir) / "slide-1.png"))],
            "source_background_warnings": ["source_background_text_removed"],
            "source_background_slide_warnings": [["source_background_table_text_skipped"]],
            "source_background_details": [{"warning_code": "source_background_detail"}],
        }

    def fake_extract_speaker_notes(_source_path: str, out_dir: str) -> list[str]:
        note_path = Path(out_dir) / "notes" / "slide-1.txt"
        note_path.parent.mkdir(parents=True, exist_ok=True)
        note_path.write_text("Extracted narration", encoding="utf-8")
        return [str(note_path)]

    monkeypatch.setattr(pptx_extract, "export_slide_images_with_metadata", fake_export_slide_images_with_metadata)
    monkeypatch.setattr(pptx_extract, "export_pptx_source_backgrounds", fake_source_backgrounds)
    monkeypatch.setattr(pptx_extract, "extract_speaker_notes", fake_extract_speaker_notes)

    slides = worker_tasks.export_project.run(str(project.id), str(source_path), False)

    assert len(slides) == 1
    assert slides[0]["source_type"] == "pptx"
    assert slides[0]["source_background_path"].replace("\\", "/").endswith("source_backgrounds/slide-1.png")
    assert slides[0]["source_background_warnings"] == [
        "source_background_text_removed",
        "source_background_table_text_skipped",
    ]
    assert slides[0]["source_render_details"] == [{"warning_code": "original_detail"}]
    assert slides[0]["source_background_details"] == [{"warning_code": "source_background_detail"}]
    assert slides[0]["whiteboard_mode"] is False
